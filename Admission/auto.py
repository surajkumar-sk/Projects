from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
import os
import sys
import json

json_data = json.loads(sys.argv[1])



wb = load_workbook('data.xlsx')
ws = wb['Form Responses 1']
max_row = ws.max_row 
column_p = 2

for k , j in json_data.items():
  ws.cell(row=max_row+1 , column=column_p).value = j
  column_p = column_p + 1


n_ROW=0
data=[]


ad = list(str(json_data['ad_num']))
an = list(str(json_data['admission_num']))
rs = list(str(json_data['r_s_number'])) if json_data['r_s_number'] else ''
m = list(str(json_data['mobile'])) if json_data['mobile'] else ''
ac = list(str(json_data['acc_num'])) if json_data['acc_num'] else ''
if_ = list(str(json_data['ifsc'])) if json_data['ifsc'] else ''
doa = json_data['doa'].split('-')
dob = json_data['dob'].split('-')
rs_date = json_data['r_s_date'].split('-') if json_data['r_s_date'] else ''

data = [str(json_data['name']) , str(json_data['fname']) , str(json_data['mname']) ,  \
  ad[0], ad[1], ad[2], ad[3], an[0] if len(an) > 0 else '' , an[1] if len(an) > 1 else '' , an[2] if len(an) > 2 else '' , an[4] if len(an) > 4 else '' , \
    ad[8] , ad[9], ad[10], ad[11] , str(dob[2]) , str(dob[1]) , str(dob[0]) , \
      str(json_data['gen']) ,  str(json_data['caste']) , str(json_data['m_tongue']) , str(json_data['rel']) , \
        str(json_data['p_l_stat']) , str(json_data['voc']) , str(json_data['l_y_a_exam']) , str(json_data['p_mark_1']) , str(json_data['p_mark_2']) , \
          str(json_data['no_d_attend']) , str(json_data['c_class']) , str(json_data['p_class']) , str(json_data['a_class']) , \
            str(json_data['med']) , str(json_data['f_lang']) , str(json_data['s_lang']) , rs[0] if len(rs) > 0 else '' , rs[1] if len(rs) > 1 else '' , rs[2] if len(rs) > 2 else '' , rs[3] if len(rs) > 3 else '' , \
              str(rs_date[2]) if rs_date else '' , str(rs_date[1]) if rs_date else ''  , str(rs_date[0]) if rs_date else ''  , str(json_data['dist']) , \
                str(json_data['mdl']) , str(json_data['mplity']) , m[0] if len(m) > 0 else '' , str(json_data['bank_n']) if json_data['bank_n'] else '' , str(json_data['branch_n']) if json_data['branch_n'] else '' ,\
                  'Nedunoor' , an[3] if len(an) > 3 else '' , str(json_data['doa']) , str(json_data['doa']) , '' , ad[4] , ad[5] , ad[6] , ad[7] , m[1] if len(m) > 0 else '' , m[2] if len(m) > 0 else '' , m[3] if len(m) > 0 else '',m[4] if len(m) > 0 else '', m[5] if len(m) > 0 else '', m[6] if len(m) > 0 else '', m[7] if len(m) > 0 else '', m[8] if len(m) > 0 else '', m[9] if len(m) > 0 else '',  \
                    ac[0] if len(ac) > 0 else '' , ac[2] if len(ac) > 2 else '' , ac[1] if len(ac) > 1 else '' , ac[3] if len(ac) > 3 else '' , ac[4] if len(ac) > 4 else '' , ac[5] if len(ac) > 5 else '' , ac[6] if len(ac) > 6 else '' , ac[7] if len(ac) > 7 else '' , ac[8] if len(ac) > 8 else '' , ac[9] if len(ac) > 9 else '' , ac[10] if len(ac) > 10 else '' , ac[11] if len(ac) > 11 else '' \
                      , ac[12] if len(ac) > 12 else '' , ac[13] if len(ac) > 13 else '' , \
                          if_[0] if len(if_) > 0 else '' , \
                            if_[2] if len(if_) > 2 else '' , if_[1] if len(if_) > 1 else '' ,  \
                              if_[3] if len(if_) > 3 else '' , if_[4] if len(if_) > 4 else '' , if_[5] if len(if_) > 5 else '' , if_[6] if len(if_) > 6 else '' , if_[7] if len(if_) > 7 else '' , if_[8] if len(if_) > 8 else '' , if_[9] if len(if_) > 9 else '' , if_[10] if len(if_) > 10 else '' , if_[11] if len(if_) > 11 else '' \
                        , if_[12] if len(if_) > 12 else '' , if_[13] if len(if_) > 13 else '']


  
path = str(os.path.dirname(os.path.abspath(__file__))) + "/Admission_form.pptx"
prs = Presentation(path)






sno=0

slide = prs.slides[0]
shapes = slide.shapes

def addingText(text_frame, text):
  if(text == 'None'):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = ''
    font = run.font
    font.name = 'Calibri (Body)'
    font.size = Pt(10.5)
    font.bold = False
    font.color.rgb = RGBColor(0,0,0)
  else:
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Calibri (Body)'
    font.size = Pt(10.5)
    font.bold = False
    font.color.rgb = RGBColor(0,0,0)


for shape in slide.shapes:
  if(shape.has_text_frame):
    addingText(shape.text_frame , data[sno])
    sno=sno+1
  


prs.save(path)
wb.save("data.xlsx") 

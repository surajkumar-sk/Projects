from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

wb = load_workbook('data.xlsx')

ws = wb['Form responses 1']

max_rows= ws.max_row - 1

path ="E:\\AUTO_PP\\excel_to_pdf_Certficates\\certificate.pptx"


def addingText(text_frame, text , font_type,color_type,font_size):
  if(text == 'None'):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = ''
    font = run.font
    font.name = font_type
    font.size = Pt(font_size)
    font.bold = True
    font.color.rgb = color_type
    # font.color.rgb = RGBColor(0,32,96)
  else:
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_type
    font.size = Pt(font_size)
    font.bold = True
    font.color.rgb = color_type


# 2 - sch_dist
# 1 - name

# prs = Presentation(path)
# slide = prs.slides[0]
# shapes = slide.shapes

# if(shapes[2].has_text_frame):
#   addingText(shapes[2].text_frame , "sch_dist" , "Calibri (Body)", RGBColor(112,48,160), 18)
# if(shapes[1].has_text_frame):
#   addingText(shapes[1].text_frame , "name" , "Calibri (Body)", RGBColor(0,0,0), 28)

score = int(str(ws.cell(row=2,column=3).value))
print(int(score) +1)
for i in range(max_rows):
  my_file = open("done.txt", "w")
  my_file.write('false')
  my_file.close()


  prs = Presentation(path)
  slide = prs.slides[0]
  shapes = slide.shapes

  score = int(str(ws.cell(row=2,column=3).value))

  if(score >= 12):

    name = str(ws.cell(row=i+2,column=4).value)
    sch_dist = str(ws.cell(row=i+2,column=5).value) + ', ' + str(ws.cell(row=i+2,column=7).value) + ' ,'

    my_file = open("data.txt", "w" , encoding='utf-8')
    my_file.write(str(ws.cell(row=i+2,column=2).value) + ' ' + str(ws.cell(row=i+2,column=7).value))
    my_file.close()

    if(shapes[2].has_text_frame):
      addingText(shapes[2].text_frame , "sch_dist" , "Calibri (Body)", RGBColor(112,48,160), 18)
    if(shapes[1].has_text_frame):
      addingText(shapes[1].text_frame , "name" , "Calibri (Body)", RGBColor(0,0,0), 28)
    
    prs.save(path)
    os.system('node converter.js')

    file1 = open("done.txt","r+")
    done = file1.read()
    file1.close()
    while(done == 'false'):
      file1 = open("done.txt","r+")
      done = file1.read()
      file1.close()
      print('happening')
  

wb.save("data.xlsx") 
prs.save(path)